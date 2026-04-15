"""
Generates docs/Pellet_Burner_Components.pptx
Run with: python3 make_pptx.py   (from inside docs/)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colour palette ────────────────────────────────────────────────────────────
BG        = RGBColor(0x1a, 0x1a, 0x2e)
CARD      = RGBColor(0x16, 0x21, 0x3e)
ACCENT    = RGBColor(0xe9, 0x45, 0x60)
TEAL      = RGBColor(0xa8, 0xda, 0xdc)
WHITE     = RGBColor(0xff, 0xff, 0xff)
GREY      = RGBColor(0xaa, 0xaa, 0xaa)
HIGHLIGHT = RGBColor(0x4c, 0xaf, 0x50)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

IMG_DIR = os.path.join(os.path.dirname(__file__), "img")

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Layout helpers ────────────────────────────────────────────────────────────
def _line_h(size):
    """Line height in inches for a given font size (pt), 1.5× leading."""
    return size * 1.5 / 72


def content_h(lines, size=14, title=None):
    """
    Compute the exact required box height (inches) for add_bullet_box content.
    Use this to set h= so nothing overflows.
    """
    h = 0.14                            # top padding
    if title:
        h += 0.38                       # title + gap
    for line in lines:
        if not line.strip():
            h += _line_h(size) * 0.5   # half-height blank line
            continue
        is_head = line.startswith("##")
        sz = (size - 1) if is_head else size
        h += _line_h(sz) + 0.06   # match the textbox height in add_bullet_box
    h += 0.12                           # bottom padding
    return h


def ch(lines, size=14, title=None, pad=0.1):
    """Return Inches(content_h(...) + pad)."""
    return Inches(content_h(lines, size, title) + pad)


# ── Drawing primitives ────────────────────────────────────────────────────────
def bg_rect(slide):
    sh = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = BG
    sh.line.fill.background()


def add_rect(slide, x, y, w, h, fill=CARD, line_color=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_bullet_box(slide, lines, x, y, w, h, size=14, title=None):
    """
    Draw a card rect and stack bullet lines inside it.
    Line heights are proportional to font size — no fixed spacing.
    h should be set using ch() so content never overflows.
    """
    add_rect(slide, x, y, w, h, fill=CARD)
    cur_y = y + Inches(0.14)

    if title:
        add_text(slide, title, x + Inches(0.15), cur_y,
                 w - Inches(0.3), Inches(0.30),
                 size=13, bold=True, color=TEAL)
        cur_y += Inches(0.38)

    for line in lines:
        if not line.strip():
            cur_y += Inches(_line_h(size) * 0.5)
            continue
        is_head = line.startswith("##")
        is_val  = line.startswith(">>")
        txt = line.lstrip("#>").strip()
        col = TEAL if is_head else (HIGHLIGHT if is_val else WHITE)
        sz  = (size - 1) if is_head else size
        lh  = Inches(_line_h(sz)) + Inches(0.06)
        prefix = "" if (is_head or is_val) else "• "
        add_text(slide, prefix + txt,
                 x + Inches(0.18), cur_y,
                 w - Inches(0.36), lh,
                 size=sz, bold=is_head, color=col)
        cur_y += lh


def place_image(slide, path, x, y, w, h, label=None):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        add_rect(slide, x, y, w, h,
                 fill=RGBColor(0x0f, 0x34, 0x60), line_color=TEAL)
        lbl = label or os.path.basename(path)
        add_text(slide, lbl, x, y + h // 2 - Inches(0.2), w, Inches(0.4),
                 size=12, color=TEAL, align=PP_ALIGN.CENTER)


def divider(slide, y):
    sh = slide.shapes.add_shape(1, Inches(0.4), y,
                                SLIDE_W - Inches(0.8), Pt(1.5))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()


def slide_title(slide, text):
    add_text(slide, text, Inches(0.4), Inches(0.15),
             Inches(12.5), Inches(0.55),
             size=28, bold=True, color=ACCENT)
    divider(slide, Inches(0.72))


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg_rect(slide)
for bar_y in (0, SLIDE_H - Inches(0.08)):
    sh = slide.shapes.add_shape(1, 0, bar_y, SLIDE_W, Inches(0.08))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

add_text(slide, "Pellet Burner Controller",
         Inches(1), Inches(2.1), Inches(11.33), Inches(1.2),
         size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "Hardware Components & Wiring Guide",
         Inches(1), Inches(3.4), Inches(11.33), Inches(0.7),
         size=26, color=TEAL, align=PP_ALIGN.CENTER)
add_text(slide,
         "ESP32  ·  SH1106 OLED  ·  DS18B20  ·  MAX6675  ·  Relays  ·  TRIAC  ·  Buttons",
         Inches(1), Inches(4.2), Inches(11.33), Inches(0.5),
         size=15, color=GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 — ESP32 DevKit V1
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg_rect(slide)
slide_title(slide, "ESP32 DevKit V1  —  Main Controller")

IMG_Y  = Inches(0.88)
IMG_X  = Inches(9.55)
IMG_W  = Inches(3.4)
IMG_H  = Inches(5.9)
place_image(slide, os.path.join(IMG_DIR, "esp32.jpg"), IMG_X, IMG_Y, IMG_W, IMG_H, "ESP32 DevKit V1")

SPECS = [
    "## Key Specifications",
    "Dual-core Xtensa LX6 @ 240 MHz",
    "520 KB SRAM  |  4 MB Flash",
    "802.11 b/g/n WiFi + Bluetooth 4.2",
    "3.3 V logic — NOT 5 V tolerant",
    "30 GPIO pins (several input-only)",
    "ADC, DAC, SPI, I2C, UART, PWM",
    "## Role in this project",
    "Main controller — FreeRTOS tasks",
    "HTTP web server for settings UI",
    "Controls relays and TRIAC output",
    "Reads DS18B20 and MAX6675 sensors",
    "Drives OLED display over I2C",
]
add_bullet_box(slide, SPECS,
               Inches(0.3), Inches(0.82), Inches(5.5), ch(SPECS))

POWER = [
    "## Power Supply",
    "USB 5V during development",
    "VIN pin: 5–12V regulated input",
    "3.3V pin: output — do NOT feed 5V",
    "Use separate 5V/3A PSU in production",
    "(relay coils draw surge current)",
]
pow_h = ch(POWER)
add_bullet_box(slide, POWER,
               Inches(5.9), Inches(0.82), Inches(3.4), pow_h)

GPIO_LINES = [
    "## GPIO Pin Assignments",
    "GPIO4   DS18B20 OneWire data",
    "GPIO5   MAX6675 SPI CS",
    "GPIO18  MAX6675 SPI SCK",
    "GPIO19  MAX6675 SPI MISO",
    "GPIO21  OLED SDA (I2C)",
    "GPIO22  OLED SCL (I2C)",
    "GPIO25  Relay — Feeder motor",
    "GPIO26  Relay — Igniter",
    "GPIO27  Relay — Water pump",
    "GPIO33  TRIAC gate trigger",
    "GPIO35  Zero-cross detect (input)",
    "GPIO13  Button: On/Off",
    "GPIO14  Button: Emergency stop",
]
gpio_y = Inches(0.82) + pow_h + Inches(0.12)
add_bullet_box(slide, GPIO_LINES, size=11,
               x=Inches(5.9), y=gpio_y,
               w=Inches(3.4), h=ch(GPIO_LINES, size=11))


# ══════════════════════════════════════════════════════════════════════════════
# Slide 3 — SH1106 OLED Display
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg_rect(slide)
slide_title(slide, "SH1106 OLED Display  —  1.3\"  128×64")

place_image(slide, os.path.join(IMG_DIR, "sh1106.jpg"),
            Inches(9.55), Inches(0.88), Inches(3.4), Inches(3.4), "SH1106 OLED")

SPECS3 = [
    "## Specifications",
    "Driver IC: SH1106  (not SSD1306)",
    "Size: 1.3 inch diagonal",
    "Resolution: 128 × 64 pixels, monochrome",
    "Interface: I2C (default address 0x3C)",
    "Operating voltage: 3.3V – 5V",
    "Library: U8g2 (SH1106 constructor)",
]
specs3_h = ch(SPECS3)
add_bullet_box(slide, SPECS3,
               Inches(0.3), Inches(0.82), Inches(5.4), specs3_h)

WIRE3 = [
    "## Wiring to ESP32",
    "VCC  →  3.3V",
    "GND  →  GND",
    "SDA  →  GPIO21",
    "SCL  →  GPIO22",
    "## Notes",
    "No pull-up resistors needed",
    "(ESP32 internal pull-ups sufficient)",
    "I2C address: 0x3C (some boards 0x3D)",
]
wire3_h = ch(WIRE3)
add_bullet_box(slide, WIRE3,
               Inches(5.85), Inches(0.82), Inches(3.5), wire3_h)

role3_y = Inches(0.82) + max(specs3_h, wire3_h) + Inches(0.14)
ROLE3 = [
    "## Role in this project",
    "At-a-glance live status without phone/browser:",
    "• Water temp (DS18B20) and flame temp (MAX6675)",
    "• Pump / Feeder / Blower state indicators",
    "• Current state machine state",
    "• Active operation mode (Standby / Auto / Timer)",
    "• Error icons (custom bitmap) on fault conditions",
    "## Important: SH1106 ≠ SSD1306",
    "Many cheap 1.3\" modules use SH1106.",
    "U8g2 constructor: U8G2_SH1106_128X64_NONAME",
]
add_bullet_box(slide, ROLE3,
               Inches(0.3), role3_y, Inches(9.0), ch(ROLE3))


# ══════════════════════════════════════════════════════════════════════════════
# Slide 4 — DS18B20 Water Temperature Sensor
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg_rect(slide)
slide_title(slide, "DS18B20  —  Water Temperature Sensor")

place_image(slide, os.path.join(IMG_DIR, "ds18b20.jpg"),
            Inches(9.55), Inches(0.88), Inches(3.4), Inches(3.5), "DS18B20")

SPECS4 = [
    "## Specifications",
    "Protocol: 1-Wire (OneWire)",
    "Range: –55°C to +125°C",
    "Accuracy: ±0.5°C (–10°C to +85°C)",
    "Resolution: 9–12 bit (configurable)",
    "Supply: 3.0V – 5.5V or parasitic (2-wire)",
    "Multiple sensors per bus (unique 64-bit ROM)",
    "Conversion: 94ms (9-bit) to 750ms (12-bit)",
]
specs4_h = ch(SPECS4)
add_bullet_box(slide, SPECS4,
               Inches(0.3), Inches(0.82), Inches(5.4), specs4_h)

WIRE4 = [
    "## Wiring — normal (3-wire) mode",
    "VCC (red)     →  3.3V",
    "GND (black)   →  GND",
    "DATA (yellow) →  GPIO4",
    "4.7 kΩ pull-up: DATA ↔ 3.3V",
    "## Wiring — parasitic (2-wire) mode",
    "VCC (red)     →  GND  (tied low)",
    "GND (black)   →  GND",
    "DATA (yellow) →  GPIO4",
    "4.7 kΩ pull-up: DATA ↔ 3.3V",
]
wire4_h = ch(WIRE4)
add_bullet_box(slide, WIRE4,
               Inches(5.85), Inches(0.82), Inches(3.5), wire4_h)

role4_y = Inches(0.82) + max(specs4_h, wire4_h) + Inches(0.14)
ROLE4 = [
    "## Role in this project",
    "Measures water / boiler temperature",
    "Used for: pump start/stop hysteresis,",
    "  power level selection (P1/P2/P3),",
    "  economy mode logic,",
    "  target temperature control loop",
    "## Recommendation",
    "Use waterproof stainless steel probe version.",
    "Route cable away from mains wiring.",
    "Use 12-bit resolution (0.0625°C steps).",
]
add_bullet_box(slide, ROLE4, size=13,
               x=Inches(0.3), y=role4_y, w=Inches(9.0), h=ch(ROLE4, size=13))


# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 — MAX6675 Thermocouple Module
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg_rect(slide)
slide_title(slide, "MAX6675  —  Flame Thermocouple Module")

place_image(slide, os.path.join(IMG_DIR, "max6675.jpg"),
            Inches(9.55), Inches(0.88), Inches(3.4), Inches(3.4), "MAX6675 Module")

SPECS5 = [
    "## Specifications",
    "Sensor type: K-type thermocouple",
    "Range: 0°C to +1023.75°C",
    "Resolution: 0.25°C (12-bit)",
    "Interface: SPI (read-only, 3-wire)",
    "Supply: 3.0V – 5.5V",
    "Cold-junction compensation: internal",
    "Open thermocouple detection (fault flag)",
    "Max SPI clock: 4.3 MHz",
]
specs5_h = ch(SPECS5)
add_bullet_box(slide, SPECS5,
               Inches(0.3), Inches(0.82), Inches(5.4), specs5_h)

WIRE5 = [
    "## Wiring to ESP32",
    "VCC       →  3.3V",
    "GND       →  GND",
    "CS        →  GPIO5",
    "SCK       →  GPIO18",
    "SO (MISO) →  GPIO19",
    "## K-type thermocouple",
    "Yellow (+)    →  T+ on module",
    "Red/White (–) →  T– on module",
    "Tip exposed to flame",
]
wire5_h = ch(WIRE5)
add_bullet_box(slide, WIRE5,
               Inches(5.85), Inches(0.82), Inches(3.5), wire5_h)

role5_y = Inches(0.82) + max(specs5_h, wire5_h) + Inches(0.14)
ROLE5 = [
    "## Role in this project",
    "Measures flame / combustion temperature",
    "Flame detection: above ~100°C = flame on",
    "Used for: ignition confirmation,",
    "  overheat protection, OLED readout",
    "## Notes",
    "Dedicated SPI bus — no sharing.",
    "Keep thermocouple wiring short and twisted.",
    "Bit-15=0 signals open thermocouple fault.",
    "Use K-type rated to at least 900°C.",
]
add_bullet_box(slide, ROLE5, size=13,
               x=Inches(0.3), y=role5_y, w=Inches(9.0), h=ch(ROLE5, size=13))


# ══════════════════════════════════════════════════════════════════════════════
# Slide 6 — Relay Modules
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg_rect(slide)
slide_title(slide, "5V Relay Modules  —  Feeder / Igniter / Pump")

place_image(slide, os.path.join(IMG_DIR, "relay.jpg"),
            Inches(9.55), Inches(0.88), Inches(3.4), Inches(3.2), "5V Relay Module")

SPECS6 = [
    "## Specifications",
    "Coil voltage: 5V DC (logic-level trigger)",
    "Trigger: active LOW  (LOW = relay ON)",
    "Contact rating: 10A @ 250V AC",
    "Isolation: optocoupler (ESP32 ↔ coil)",
    "Flyback diode: included on module",
    "Indicator LED per channel",
]
specs6_h = ch(SPECS6)
add_bullet_box(slide, SPECS6,
               Inches(0.3), Inches(0.82), Inches(5.4), specs6_h)

relay_row_y = Inches(0.82) + specs6_h + Inches(0.14)

for i, (name, gpio, desc) in enumerate([
    ("Feeder Motor", "GPIO25", "Pellet auger / conveyor"),
    ("Igniter",      "GPIO26", "Electric heating element"),
    ("Water Pump",   "GPIO27", "Circulation pump"),
]):
    RL = [
        f"## {name}",
        f"IN  →  {gpio}",
        "VCC →  5V",
        "GND →  GND",
        "COM →  AC Live",
        "NO  →  Load",
        desc,
    ]
    bx = Inches(0.3) + i * Inches(3.1)
    add_bullet_box(slide, RL, size=13,
                   x=bx, y=relay_row_y, w=Inches(3.0), h=ch(RL, size=13))

SAFETY6 = [
    "## Active LOW wiring note",
    "ESP32 GPIO init HIGH at boot →",
    "relay stays OFF during reset/startup",
    "LOW = relay ON,  HIGH = relay OFF",
    "## Safety",
    "All relays default OFF on power-up",
    "Emergency stop: all GPIOs → HIGH",
]
add_bullet_box(slide, SAFETY6, size=13,
               x=Inches(9.55), y=relay_row_y,
               w=Inches(3.4), h=ch(SAFETY6, size=13))


# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 — TRIAC Blower Control
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg_rect(slide)
slide_title(slide, "TRIAC Dimmer  —  Blower Motor Speed Control")

place_image(slide, os.path.join(IMG_DIR, "triac.jpg"),
            Inches(9.55), Inches(0.88), Inches(3.4), Inches(3.8), "TRIAC / AC Dimmer")

METHOD = [
    "## Method: Phase-angle control",
    "The AC waveform is chopped each half-cycle.",
    "TRIAC fires later in the cycle → less power.",
    "Zero-cross detector tells ESP32 exactly when",
    "the AC wave crosses 0V.",
    "ESP32 waits N µs, then pulses the TRIAC gate",
    "→ conducts for the rest of the half-cycle.",
    "Larger delay = lower fan speed.",
    "Delay = 0 = full power.",
    "## Typical components",
    "BT136 or BTA16 TRIAC (TO-220)",
    "MOC3021 optocoupler (gate isolation)",
    "4N35 optocoupler (zero-cross isolation)",
    "Snubber: 100Ω + 100nF/400V across TRIAC",
    "Gate resistor: 330Ω before MOC3021",
]
add_bullet_box(slide, METHOD, size=13,
               x=Inches(0.3), y=Inches(0.82),
               w=Inches(5.4), h=ch(METHOD, size=13))

WIRE7 = [
    "## Zero-cross input",
    "GPIO35  ←  Zero-cross signal",
    "(GPIO35 is input-only on ESP32)",
    "Signal: HIGH normally, drops LOW at",
    "each zero crossing of the AC wave",
    "Trigger ISR on falling edge",
    "## TRIAC gate output",
    "GPIO33  →  330Ω  →  MOC3021 LED(+)",
    "MOC3021 output  →  TRIAC gate",
    "## Voltage levels",
    "ESP32 side: 3.3V logic (optocoupled)",
    "AC side: 230V AC — HIGH VOLTAGE",
    "Keep AC and DC wiring separated",
    "Fuse the AC circuit (5–10A slow blow)",
]
add_bullet_box(slide, WIRE7, size=13,
               x=Inches(5.85), y=Inches(0.82),
               w=Inches(3.5), h=ch(WIRE7, size=13))


# ══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Physical Buttons
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg_rect(slide)
slide_title(slide, "Physical Buttons  —  On/Off & Emergency Stop")

place_image(slide, os.path.join(IMG_DIR, "button.jpg"),
            Inches(9.55), Inches(0.88), Inches(3.4), Inches(3.0), "Push Button")

TYPE8 = [
    "## Button type",
    "Momentary normally-open (NO) push button",
    "Panel-mount, rated for frequent use",
    "On/Off: latching toggle preferred",
    "E-Stop: momentary (NO) only",
    "## Wiring  (active LOW, internal pull-up)",
    "One terminal  →  GPIO pin",
    "Other terminal  →  GND",
    "ESP32 internal pull-up enabled in firmware",
    "No external resistors needed",
    "Button pressed  →  GPIO reads LOW",
    "Button released  →  GPIO reads HIGH",
]
type8_h = ch(TYPE8)
add_bullet_box(slide, TYPE8,
               Inches(0.3), Inches(0.82), Inches(5.4), type8_h)

BTN8 = [
    "## Button 1 — On / Off",
    "GPIO13  ←  button  ←  GND",
    "Toggles Standby ↔ previous mode",
    "(Automatic or Timer)",
    "",
    "## Button 2 — Emergency Stop",
    "GPIO14  ←  button  ←  GND",
    "Immediate shutdown — ALL outputs OFF",
    "State machine → STANDBY instantly",
    "Does NOT write op_mode to NVS flash",
]
btn8_h = ch(BTN8)
add_bullet_box(slide, BTN8,
               Inches(5.85), Inches(0.82), Inches(3.5), btn8_h)

deb_y = Inches(0.82) + max(type8_h, btn8_h) + Inches(0.14)
DEB8 = [
    "## Debounce strategy",
    "Hardware: 100nF cap from GPIO pin to GND (optional).",
    "Software: ignore transitions for 50ms after first edge.",
    "E-Stop: no debounce delay — act on first LOW edge immediately.",
]
add_bullet_box(slide, DEB8, size=13,
               x=Inches(0.3), y=deb_y, w=Inches(9.0), h=ch(DEB8, size=13))


# ══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Complete Wiring Summary
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg_rect(slide)
slide_title(slide, "Complete Wiring Summary")

ROW1_Y = Inches(0.82)
COL1, COL2, COL3 = Inches(0.3), Inches(4.55), Inches(8.8)

I2C = [
    "## I2C Bus  (GPIO21 SDA, GPIO22 SCL)",
    "SH1106 OLED  —  addr 0x3C",
    "VCC → 3.3V   GND → GND",
    "SDA → GPIO21   SCL → GPIO22",
    "## OneWire Bus  (GPIO4)",
    "DS18B20 water temp sensor",
    "DATA → GPIO4  (4.7kΩ to 3.3V)",
    "VCC → 3.3V   GND → GND",
]
SPI = [
    "## SPI Bus  (dedicated)",
    "MAX6675 thermocouple module",
    "CS   → GPIO5",
    "SCK  → GPIO18",
    "MISO → GPIO19",
    "VCC  → 3.3V   GND → GND",
    "## Zero-cross detect",
    "Signal → GPIO35 (input only)",
    "3.3V logic, falling edge ISR",
]
RELAYS = [
    "## Relay outputs  (active LOW)",
    "GPIO25 → Relay IN — Feeder motor",
    "GPIO26 → Relay IN — Igniter",
    "GPIO27 → Relay IN — Pump",
    "Relay VCC → 5V   GND → GND",
    "## TRIAC gate",
    "GPIO33 → 330Ω → MOC3021 → TRIAC",
]

row1_h = ch(I2C)
add_bullet_box(slide, I2C,   COL1, ROW1_Y, Inches(4.1), row1_h)
add_bullet_box(slide, SPI,   COL2, ROW1_Y, Inches(4.1), ch(SPI))
add_bullet_box(slide, RELAYS, COL3, ROW1_Y, Inches(4.2), ch(RELAYS))

ROW2_Y = ROW1_Y + max(ch(I2C), ch(SPI), ch(RELAYS)) + Inches(0.14)

BTNS = [
    "## Buttons  (active LOW, internal pull-up)",
    "GPIO13 ← On/Off button ← GND",
    "GPIO14 ← Emergency Stop button ← GND",
]
PWR = [
    "## Power rails",
    "3.3V: ESP32 reg, OLED, DS18B20, MAX6675",
    "5V: Relay module VCC",
    "AC 230V: TRIAC load side — isolated",
]
SAFE = [
    "## Safety reminders",
    "Never connect 5V to ESP32 3.3V rail",
    "Optocouplers required on TRIAC side",
    "Fuse the AC circuit (5–10A slow blow)",
]
add_bullet_box(slide, BTNS, COL1, ROW2_Y, Inches(4.1), ch(BTNS))
add_bullet_box(slide, PWR,  COL2, ROW2_Y, Inches(4.1), ch(PWR))
add_bullet_box(slide, SAFE, COL3, ROW2_Y, Inches(4.2), ch(SAFE))


# ══════════════════════════════════════════════════════════════════════════════
# Slide 10 — Pin Reference Table
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg_rect(slide)
slide_title(slide, "Pin Reference Table")

headers = ["Component", "GPIO", "Protocol", "Voltage", "Notes"]
rows = [
    ["SH1106 OLED — SDA",   "GPIO21", "I2C",    "3.3 V", "I2C address 0x3C"],
    ["SH1106 OLED — SCL",   "GPIO22", "I2C",    "3.3 V", "400 kHz fast-mode"],
    ["DS18B20 DATA",        "GPIO4",  "1-Wire", "3.3 V", "4.7 kΩ pull-up to 3.3 V"],
    ["MAX6675 CS",          "GPIO5",  "SPI",    "3.3 V", "dedicated SPI bus"],
    ["MAX6675 SCK",         "GPIO18", "SPI",    "3.3 V", "max 4.3 MHz"],
    ["MAX6675 MISO (SO)",   "GPIO19", "SPI",    "3.3 V", "read-only from module"],
    ["Relay — Feeder IN",   "GPIO25", "GPIO",   "3.3 V", "active LOW, relay VCC = 5 V"],
    ["Relay — Igniter IN",  "GPIO26", "GPIO",   "3.3 V", "active LOW, relay VCC = 5 V"],
    ["Relay — Pump IN",     "GPIO27", "GPIO",   "3.3 V", "active LOW, relay VCC = 5 V"],
    ["TRIAC gate",          "GPIO33", "GPIO",   "3.3 V", "330 Ω → MOC3021 → TRIAC gate"],
    ["Zero-cross detect",   "GPIO35", "GPIO",   "3.3 V", "input-only pin, falling edge ISR"],
    ["Button — On/Off",     "GPIO13", "GPIO",   "3.3 V", "internal pull-up, active LOW"],
    ["Button — E-Stop",     "GPIO14", "GPIO",   "3.3 V", "internal pull-up, active LOW"],
]

NUM_ROWS = 1 + len(rows)
NUM_COLS = len(headers)
TBL_LEFT = Inches(0.3)
TBL_TOP  = Inches(0.88)
TBL_W    = SLIDE_W - Inches(0.6)
TBL_H    = SLIDE_H - Inches(1.05)

tbl_shape = slide.shapes.add_table(
    NUM_ROWS, NUM_COLS, TBL_LEFT, TBL_TOP, int(TBL_W), int(TBL_H))
tbl = tbl_shape.table

col_fracs = [0.26, 0.10, 0.10, 0.09, 0.45]
for ci, frac in enumerate(col_fracs):
    tbl.columns[ci].width = int(TBL_W * frac)

HDR_H  = Inches(0.48)
DATA_H = int((TBL_H - HDR_H) / len(rows))
tbl.rows[0].height = int(HDR_H)
for ri in range(1, NUM_ROWS):
    tbl.rows[ri].height = DATA_H


def set_cell(tbl, ri, ci, text, bold=False, font_size=13,
             fg=WHITE, bg_color=None, align=PP_ALIGN.LEFT):
    cell = tbl.cell(ri, ci)
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = fg
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    sf   = etree.SubElement(tcPr, qn('a:solidFill'))
    clr  = etree.SubElement(sf, qn('a:srgbClr'))
    c    = bg_color if bg_color else BG
    clr.set('val', f'{int(c[0]):02X}{int(c[1]):02X}{int(c[2]):02X}')
    for side in ('lnL', 'lnR', 'lnT', 'lnB'):
        ln = etree.SubElement(tcPr, qn(f'a:{side}'))
        etree.SubElement(ln, qn('a:noFill'))


for ci, hdr in enumerate(headers):
    set_cell(tbl, 0, ci, hdr, bold=True, font_size=14,
             fg=WHITE, bg_color=ACCENT, align=PP_ALIGN.CENTER)

ROW_EVEN = RGBColor(0x0f, 0x34, 0x60)
for ri, row in enumerate(rows):
    bg = ROW_EVEN if ri % 2 == 0 else CARD
    for ci, val in enumerate(row):
        fg = TEAL if ci == 1 else (GREY if ci == 4 else WHITE)
        set_cell(tbl, ri + 1, ci, val, font_size=13, fg=fg, bg_color=bg)


# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "Pellet_Burner_Components.pptx")
prs.save(out)
print(f"Saved: {out}")
